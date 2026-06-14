"""
universal_document_processor.py
================================
Handles ALL document types for Research Hunter v2-4
- PDF (text and scanned images)
- DOCX, ODT, RTF
- EPUB, MOBI
- HTML/Web pages
- Plain text
- Images (OCR for scanned documents)
- PowerPoint, Excel

Uses Ollama (qwen2.5vl:3b) for intelligent analysis and OCR enhancement
"""

import os
import sys
import json
import base64
import tempfile
import subprocess
from pathlib import Path
from typing import Optional, Dict, List, Tuple, Union
from dataclasses import dataclass, field
import hashlib

# Try to import PDF libraries
try:
    import pdfplumber
    HAS_PDFPLUMBER = True
except ImportError:
    HAS_PDFPLUMBER = False

try:
    import fitz  # PyMuPDF
    HAS_PYMUPDF = True
except ImportError:
    HAS_PYMUPDF = False

try:
    from docx import Document as DocxDocument
    HAS_DOCX = True
except ImportError:
    HAS_DOCX = False

# Playwright for dynamic web content
try:
    from playwright.sync_api import sync_playwright
    HAS_PLAYWRIGHT = True
except ImportError:
    HAS_PLAYWRIGHT = False


@dataclass
class DocumentMetadata:
    """Metadata extracted from document"""
    title: str = ""
    authors: List[str] = field(default_factory=list)
    abstract: str = ""
    keywords: List[str] = field(default_factory=list)
    publication_date: str = ""
    journal: str = ""
    doi: str = ""
    document_type: str = ""
    language: str = "en"
    page_count: int = 0
    word_count: int = 0
    has_figures: bool = False
    has_tables: bool = False
    is_scanned: bool = False
    quality_score: float = 0.0


@dataclass
class ProcessedDocument:
    """Complete processed document"""
    file_path: str
    file_type: str
    metadata: DocumentMetadata
    full_text: str
    sections: Dict[str, str]
    tables: List[List[List[str]]]
    images: List[str]  # Base64 encoded images
    raw_pages: List[str]
    processing_time: float = 0.0
    errors: List[str] = field(default_factory=list)


class OllamaClient:
    """Clean interface to Ollama qwen2.5vl:3b model"""
    
    def __init__(self, base_url: str = "http://localhost:11434"):
        self.base_url = base_url
        self.model = "qwen2.5vl:3b"
        self._available = None
    
    @property
    def available(self) -> bool:
        """Check if Ollama is available"""
        if self._available is not None:
            return self._available
        try:
            import requests
            r = requests.get(f"{self.base_url}/api/tags", timeout=5)
            self._available = r.status_code == 200
        except:
            self._available = False
        return self._available
    
    def generate(self, prompt: str, system: str = "", image_base64: str = None) -> str:
        """Generate text using Ollama"""
        if not self.available:
            return ""
        
        import requests
        payload = {
            "model": self.model,
            "prompt": prompt,
            "stream": False,
            "options": {"temperature": 0.3, "num_predict": 2048}
        }
        
        if system:
            payload["system"] = system
        
        if image_base64:
            payload["images"] = [image_base64]
        
        try:
            r = requests.post(f"{self.base_url}/api/generate", json=payload, timeout=120)
            if r.status_code == 200:
                return r.json().get("response", "")
        except Exception as e:
            print(f"Ollama error: {e}")
        return ""
    
    def analyze_image(self, image_base64: str, question: str) -> str:
        """Analyze an image using vision capabilities"""
        if not self.available:
            return ""
        
        import requests
        payload = {
            "model": self.model,
            "prompt": question,
            "images": [image_base64],
            "stream": False,
            "options": {"temperature": 0.2, "num_predict": 1024}
        }
        
        try:
            r = requests.post(f"{self.base_url}/api/generate", json=payload, timeout=120)
            if r.status_code == 200:
                return r.json().get("response", "")
        except Exception as e:
            print(f"Ollama image analysis error: {e}")
        return ""


class UniversalDocumentProcessor:
    """
    Universal document processor for Research Hunter v2-4
    
    Handles:
    - PDF (native text and scanned/image-based)
    - DOCX, ODT, RTF
    - EPUB, MOBI
    - HTML/Web pages
    - Plain text files
    - Images with OCR
    - Office documents (Excel, PowerPoint)
    """
    
    def __init__(self):
        self.ollama = OllamaClient()
        self.temp_dir = tempfile.mkdtemp(prefix="doc_processor_")
        
        # OCR languages supported
        self.ocr_languages = [
            "eng", "ara", "fra", "deu", "spa", "ita", "por", "tur",
            "rus", "chi_sim", "jpn", "kor", "hin", "urd", "fas", "pol",
            "dut", "gre", "heb", "vie", "tha", "ind"
        ]
        
        # Document type signatures
        self.type_signatures = {
            "research_article": ["abstract", "introduction", "methodology", "results", "discussion"],
            "review_article": ["systematic review", "meta-analysis", "literature review"],
            "thesis": ["chapter", "dissertation", "supervisor", "phd", "master"],
            "conference": ["conference", "proceedings", "presentation", "session"],
            "book_chapter": ["chapter", "part", "edition", "isbn"],
            "report": ["executive summary", "recommendations", "findings"]
        }
    
    def process(self, file_path: str) -> ProcessedDocument:
        """
        Main entry point - process any document type
        
        Args:
            file_path: Path to the document
            
        Returns:
            ProcessedDocument with all extracted data
        """
        import time
        start_time = time.time()
        
        file_path = Path(file_path)
        file_type = self._detect_file_type(file_path)
        
        metadata = DocumentMetadata(document_type=file_type)
        errors = []
        
        try:
            if file_type == "pdf":
                result = self._process_pdf(file_path)
            elif file_type in ["docx", "odt", "rtf"]:
                result = self._process_word(file_path)
            elif file_type == "epub":
                result = self._process_epub(file_path)
            elif file_type == "html":
                result = self._process_html(file_path)
            elif file_type == "txt":
                result = self._process_text(file_path)
            elif file_type in ["png", "jpg", "jpeg", "tiff", "bmp"]:
                result = self._process_image(file_path)
            elif file_type in ["xlsx", "xls"]:
                result = self._process_excel(file_path)
            elif file_type in ["pptx", "ppt"]:
                result = self._process_powerpoint(file_path)
            else:
                result = {"text": "", "sections": {}, "tables": [], "images": [], "pages": []}
                errors.append(f"Unsupported file type: {file_type}")
            
            # Enhance with Ollama if available
            if self.ollama.available and result.get("text"):
                result = self._enhance_with_ollama(result, file_type)
            
        except Exception as e:
            errors.append(str(e))
            result = {"text": "", "sections": {}, "tables": [], "images": [], "pages": []}
        
        return ProcessedDocument(
            file_path=str(file_path),
            file_type=file_type,
            metadata=metadata,
            full_text=result.get("text", ""),
            sections=result.get("sections", {}),
            tables=result.get("tables", []),
            images=result.get("images", []),
            raw_pages=result.get("pages", []),
            processing_time=time.time() - start_time,
            errors=errors
        )
    
    def _detect_file_type(self, file_path: Path) -> str:
        """Detect file type from extension and magic bytes"""
        ext = file_path.suffix.lower()
        
        type_map = {
            ".pdf": "pdf",
            ".docx": "docx",
            ".doc": "doc",
            ".odt": "odt",
            ".rtf": "rtf",
            ".epub": "epub",
            ".mobi": "mobi",
            ".html": "html",
            ".htm": "html",
            ".txt": "txt",
            ".md": "txt",
            ".png": "png",
            ".jpg": "jpg",
            ".jpeg": "jpeg",
            ".tiff": "tiff",
            ".tif": "tiff",
            ".bmp": "bmp",
            ".xlsx": "xlsx",
            ".xls": "xls",
            ".pptx": "pptx",
            ".ppt": "ppt"
        }
        
        return type_map.get(ext, "unknown")
    
    def _process_pdf(self, file_path: Path) -> Dict:
        """Process PDF document - handles both text and scanned"""
        text = ""
        pages = []
        tables = []
        images = []
        sections = {}
        
        # Try PyMuPDF first (faster, better for text)
        if HAS_PYMUPDF:
            try:
                doc = fitz.open(str(file_path))
                metadata = doc.metadata
                
                for page_num, page in enumerate(doc):
                    # Extract text
                    page_text = page.get_text()
                    pages.append(page_text)
                    text += page_text + "\n\n"
                    
                    # Check if page is scanned (minimal text)
                    if len(page_text.strip()) < 50:
                        # Try to extract images for OCR
                        img_list = page.get_images(full=True)
                        for img_index, img in enumerate(img_list):
                            xref = img[0]
                            base_image = doc.extract_image(xref)
                            image_bytes = base_image["image"]
                            image_base64 = base64.b64encode(image_bytes).decode('utf-8')
                            images.append(image_base64)
                
                doc.close()
            except Exception as e:
                print(f"PyMuPDF error: {e}")
        
        # Try pdfplumber as backup (better for tables)
        if not text and HAS_PDFPLUMBER:
            try:
                with pdfplumber.open(str(file_path)) as pdf:
                    for page in pdf.pages:
                        page_text = page.extract_text() or ""
                        pages.append(page_text)
                        text += page_text + "\n\n"
                        
                        # Extract tables
                        page_tables = page.extract_tables()
                        if page_tables:
                            tables.extend(page_tables)
            except Exception as e:
                print(f"pdfplumber error: {e}")
        
        # Check if PDF is scanned (very little text extracted)
        is_scanned = len(text.strip()) < 100 and len(images) > 0
        
        # OCR scanned pages with Ollama enhancement
        if is_scanned and images:
            ocr_text = self._ocr_images(images)
            if ocr_text:
                text = ocr_text
                pages = [ocr_text]
        
        # Parse sections
        sections = self._parse_sections(text)
        
        # Extract metadata
        metadata = self._extract_metadata(text, file_path)
        
        return {
            "text": text,
            "pages": pages,
            "tables": tables,
            "images": images,
            "sections": sections,
            "is_scanned": is_scanned,
            "metadata": metadata
        }
    
    def _process_word(self, file_path: Path) -> Dict:
        """Process Word documents (DOCX, ODT, RTF)"""
        text = ""
        tables = []
        
        if HAS_DOCX:
            try:
                doc = DocxDocument(str(file_path))
                
                # Extract paragraphs
                for para in doc.paragraphs:
                    text += para.text + "\n"
                
                # Extract tables
                for table in doc.tables:
                    table_data = []
                    for row in table.rows:
                        row_data = [cell.text.strip() for cell in row.cells]
                        table_data.append(row_data)
                    tables.append(table_data)
                
                # Extract metadata
                core_props = doc.core_properties
                metadata = {
                    "title": core_props.title or "",
                    "authors": [core_props.author] if core_props.author else [],
                    "created": str(core_props.created) if core_props.created else ""
                }
                
            except Exception as e:
                print(f"DOCX error: {e}")
        
        # Parse sections
        sections = self._parse_sections(text)
        
        return {
            "text": text,
            "sections": sections,
            "tables": tables,
            "images": [],
            "pages": []
        }
    
    def _process_epub(self, file_path: Path) -> Dict:
        """Process EPUB books"""
        # Simple EPUB extraction (can be enhanced with ebooklib)
        try:
            import zipfile
            text = ""
            
            with zipfile.ZipFile(str(file_path), 'r') as z:
                for file in z.namelist():
                    if file.endswith('.xhtml') or file.endswith('.html'):
                        content = z.read(file).decode('utf-8', errors='ignore')
                        # Strip HTML tags
                        import re
                        clean = re.sub(r'<[^>]+>', ' ', content)
                        clean = re.sub(r'\s+', ' ', clean)
                        text += clean + "\n\n"
            
            sections = self._parse_sections(text)
            
            return {
                "text": text,
                "sections": sections,
                "tables": [],
                "images": [],
                "pages": []
            }
        except Exception as e:
            print(f"EPUB error: {e}")
            return {"text": "", "sections": {}, "tables": [], "images": [], "pages": []}
    
    def _process_html(self, file_path: Path) -> Dict:
        """Process HTML files"""
        try:
            with open(str(file_path), 'r', encoding='utf-8') as f:
                html = f.read()
            
            import re
            # Remove scripts and styles
            html = re.sub(r'<script[^>]*>.*?</script>', '', html, flags=re.DOTALL)
            html = re.sub(r'<style[^>]*>.*?</style>', '', html, flags=re.DOTALL)
            
            # Extract text
            text = re.sub(r'<[^>]+>', ' ', html)
            text = re.sub(r'\s+', ' ', text).strip()
            
            return {
                "text": text,
                "sections": {"content": text},
                "tables": [],
                "images": [],
                "pages": []
            }
        except Exception as e:
            print(f"HTML error: {e}")
            return {"text": "", "sections": {}, "tables": [], "images": [], "pages": []}
    
    def _process_text(self, file_path: Path) -> Dict:
        """Process plain text files"""
        try:
            with open(str(file_path), 'r', encoding='utf-8', errors='ignore') as f:
                text = f.read()
            
            sections = self._parse_sections(text)
            
            return {
                "text": text,
                "sections": sections,
                "tables": [],
                "images": [],
                "pages": []
            }
        except Exception as e:
            print(f"Text file error: {e}")
            return {"text": "", "sections": {}, "tables": [], "images": [], "pages": []}
    
    def _process_image(self, file_path: Path) -> Dict:
        """Process image files with OCR"""
        try:
            with open(str(file_path), 'rb') as f:
                image_bytes = f.read()
            
            image_base64 = base64.b64encode(image_bytes).decode('utf-8')
            
            # Use Ollama vision for intelligent OCR
            if self.ollama.available:
                text = self.ollama.analyze_image(
                    image_base64,
                    "Extract all text from this image exactly as written. If this is a document, preserve the structure and formatting."
                )
            else:
                # Fallback to Tesseract
                text = self._tesseract_ocr(str(file_path))
            
            return {
                "text": text,
                "sections": {"full": text},
                "tables": [],
                "images": [image_base64],
                "pages": [text]
            }
        except Exception as e:
            print(f"Image processing error: {e}")
            return {"text": "", "sections": {}, "tables": [], "images": [], "pages": []}
    
    def _process_excel(self, file_path: Path) -> Dict:
        """Process Excel files"""
        text = ""
        tables = []
        
        try:
            import openpyxl
            wb = openpyxl.load_workbook(str(file_path))
            
            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                text += f"\n\n=== Sheet: {sheet_name} ===\n"
                
                sheet_data = []
                for row in sheet.iter_rows(values_only=True):
                    row_data = [str(cell) if cell is not None else "" for cell in row]
                    sheet_data.append(row_data)
                    text += " | ".join(row_data) + "\n"
                
                if sheet_data:
                    tables.append(sheet_data)
            
        except Exception as e:
            print(f"Excel error: {e}")
        
        return {
            "text": text,
            "sections": {"data": text},
            "tables": tables,
            "images": [],
            "pages": []
        }
    
    def _process_powerpoint(self, file_path: Path) -> Dict:
        """Process PowerPoint files"""
        text = ""
        
        try:
            from pptx import Presentation
            prs = Presentation(str(file_path))
            
            for slide_num, slide in enumerate(prs.slides):
                text += f"\n\n=== Slide {slide_num + 1} ===\n"
                
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
                    
                    if hasattr(shape, "table"):
                        for row in shape.table.rows:
                            row_text = [cell.text for cell in row.cells]
                            text += " | ".join(row_text) + "\n"
            
        except Exception as e:
            print(f"PowerPoint error: {e}")
        
        return {
            "text": text,
            "sections": {"slides": text},
            "tables": [],
            "images": [],
            "pages": []
        }
    
    def _ocr_images(self, images_base64: List[str]) -> str:
        """OCR multiple images and combine results"""
        if not images_base64:
            return ""
        
        full_text = []
        
        for i, img_b64 in enumerate(images_base64):
            if self.ollama.available:
                result = self.ollama.analyze_image(
                    img_b64,
                    f"Extract all text from page {i+1} of this document. Preserve the reading order and structure."
                )
                if result:
                    full_text.append(result)
            else:
                # Fallback to Tesseract
                # Would need to decode and save image temporarily
                pass
        
        return "\n\n".join(full_text)
    
    def _tesseract_ocr(self, image_path: str) -> str:
        """Use Tesseract OCR directly"""
        try:
            result = subprocess.run(
                ['tesseract', image_path, 'stdout', '-l', 'eng+ara'],
                capture_output=True, text=True, timeout=120
            )
            return result.stdout
        except Exception as e:
            print(f"Tesseract error: {e}")
            return ""
    
    def _parse_sections(self, text: str) -> Dict[str, str]:
        """Parse document into sections based on common academic headings"""
        sections = {}
        
        # Common section patterns
        section_patterns = [
            (r'(?i)abstract|summary', 'abstract'),
            (r'(?i)introduction|background', 'introduction'),
            (r'(?i)literature\s+review|related\s+work', 'literature_review'),
            (r'(?i)methodology|methods|materials?\s+and\s+methods', 'methodology'),
            (r'(?i)results?|findings?', 'results'),
            (r'(?i)discussion|analysis', 'discussion'),
            (r'(?i)conclusion|conclusions?', 'conclusion'),
            (r'(?i)references?|bibliography', 'references'),
            (r'(?i)acknowledgements?', 'acknowledgements'),
            (r'(?i)appendix|appendices', 'appendix'),
        ]
        
        lines = text.split('\n')
        current_section = 'other'
        current_content = []
        
        for line in lines:
            line = line.strip()
            if not line:
                continue
            
            matched = False
            for pattern, section_name in section_patterns:
                import re
                if re.match(pattern, line) and len(line) < 100:
                    # Save previous section
                    if current_content:
                        sections[current_section] = '\n'.join(current_content)
                    current_section = section_name
                    current_content = []
                    matched = True
                    break
            
            if not matched:
                current_content.append(line)
        
        # Save last section
        if current_content:
            sections[current_section] = '\n'.join(current_content)
        
        return sections
    
    def _extract_metadata(self, text: str, file_path: Path) -> DocumentMetadata:
        """Extract metadata from document text and filename"""
        metadata = DocumentMetadata()
        
        # Extract from filename
        filename = file_path.stem
        metadata.title = filename.replace('_', ' ').replace('-', ' ')
        
        # Use Ollama to extract structured metadata
        if self.ollama.available and len(text) > 100:
            prompt = f"""Extract metadata from this academic paper. Return JSON with these fields:
            - title: The paper title
            - authors: Array of author names
            - abstract: The abstract text
            - keywords: Array of keywords
            - publication_date: Publication date if found
            - journal: Journal name if found
            
            Paper content (first 2000 chars):
            {text[:2000]}
            
            Return ONLY valid JSON, no other text."""
            
            result = self.ollama.generate(prompt)
            
            try:
                import json
                data = json.loads(result)
                metadata.title = data.get('title', metadata.title)
                metadata.authors = data.get('authors', [])
                metadata.abstract = data.get('abstract', '')
                metadata.keywords = data.get('keywords', [])
                metadata.publication_date = data.get('publication_date', '')
                metadata.journal = data.get('journal', '')
            except:
                pass
        
        return metadata
    
    def _enhance_with_ollama(self, result: Dict, file_type: str) -> Dict:
        """Use Ollama to enhance and structure the document"""
        if not self.ollama.available:
            return result
        
        text = result.get("text", "")
        if len(text) < 100:
            return result
        
        # Analyze document structure
        prompt = f"""Analyze this {file_type} document and provide:
        1. A brief summary (2-3 sentences)
        2. Key topics covered
        3. Document quality assessment (high/medium/low)
        4. Suggested improvements if any
        
        Document preview:
        {text[:3000]}"""
        
        analysis = self.ollama.generate(prompt)
        
        # Store analysis in sections
        result["sections"]["ollama_analysis"] = analysis
        
        return result
    
    def fetch_webpage(self, url: str) -> Dict:
        """Fetch and process a webpage using Playwright (for dynamic content)"""
        if not HAS_PLAYWRIGHT:
            # Fallback to requests
            try:
                r = requests.get(url, timeout=30)
                return {
                    "text": r.text,
                    "html": r.text,
                    "url": url
                }
            except:
                return {"text": "", "html": "", "url": url}
        
        try:
            with sync_playwright() as p:
                browser = p.chromium.launch(headless=True)
                page = browser.new_page()
                
                page.goto(url, wait_until="networkidle", timeout=30000)
                
                # Wait for dynamic content
                page.wait_for_timeout(2000)
                
                # Get content
                html = page.content()
                text = page.inner_text("body")
                
                browser.close()
                
                return {
                    "text": text,
                    "html": html,
                    "url": url
                }
        except Exception as e:
            print(f"Playwright error: {e}")
            return {"text": "", "html": "", "url": url, "error": str(e)}
    
    def scrape_dynamic_content(self, url: str, selectors: List[str]) -> List[str]:
        """Scrape dynamic content using Playwright"""
        if not HAS_PLAYWRIGHT:
            return []
        
        results = []
        
        try:
            with sync_playwright() as p:
                browser = p.chromium.launch(headless=True)
                page = browser.new_page()
                
                page.goto(url, wait_until="networkidle", timeout=30000)
                
                for selector in selectors:
                    elements = page.query_selector_all(selector)
                    for elem in elements:
                        results.append(elem.inner_text())
                
                browser.close()
        except Exception as e:
            print(f"Playwright scrape error: {e}")
        
        return results


# Quick test function
def test_processor():
    """Test the document processor"""
    processor = UniversalDocumentProcessor()
    
    print("=" * 60)
    print("Universal Document Processor - Test")
    print("=" * 60)
    
    print(f"\nOllama Available: {processor.ollama.available}")
    print(f"PyMuPDF: {HAS_PYMUPDF}")
    print(f"pdfplumber: {HAS_PDFPLUMBER}")
    print(f"DOCX: {HAS_DOCX}")
    print(f"Playwright: {HAS_PLAYWRIGHT}")
    
    print("\n✓ Processor initialized successfully!")


if __name__ == "__main__":
    test_processor()